"""
Create a PowerPoint presentation for the Fraud Detection project
using visualizations from the outputs directory.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

OUTPUT_DIR = 'D:/Final Year/App Domains 3/datasets/outputs'

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullet_points=None, image_path=None, image_width=None):
    """Add a content slide with optional image"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1),
                                   Inches(12.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 102, 204)
    line.line.fill.background()

    if image_path and os.path.exists(image_path):
        if bullet_points:
            # Image on right, bullets on left
            img_width = image_width or Inches(7)
            slide.shapes.add_picture(image_path, Inches(5.8), Inches(1.4), width=img_width)

            # Bullets on left
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5.5))
            tf = text_box.text_frame
            tf.word_wrap = True

            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = "• " + point
                p.font.size = Pt(18)
                p.space_after = Pt(12)
        else:
            # Just image, centered
            img_width = image_width or Inches(11)
            img_left = (Inches(13.333) - img_width) / 2
            slide.shapes.add_picture(image_path, img_left, Inches(1.4), width=img_width)
    elif bullet_points:
        # Just bullets
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = text_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(22)
            p.space_after = Pt(14)

    return slide

def add_two_image_slide(prs, title, image1_path, image2_path, caption1="", caption2=""):
    """Add a slide with two images side by side"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Images
    if os.path.exists(image1_path):
        slide.shapes.add_picture(image1_path, Inches(0.3), Inches(1.3), width=Inches(6.2))
    if os.path.exists(image2_path):
        slide.shapes.add_picture(image2_path, Inches(6.8), Inches(1.3), width=Inches(6.2))

    return slide

# =========================================================================
# CREATE SLIDES
# =========================================================================

print("Creating presentation...")

# Slide 1: Title
add_title_slide(prs,
    "Fraud Detection in Financial Transactions\nUsing Machine Learning",
    "Lee Hickey | Emma Reen | Sarah O'Hanlon\nDublin City University")

# Slide 2: Agenda
add_content_slide(prs, "Agenda", [
    "Introduction & Research Questions",
    "Dataset Overview & Exploratory Analysis",
    "Critical Finding: Data Leakage Prevention",
    "RQ1: Algorithm Comparison",
    "RQ2: Feature Importance Analysis",
    "RQ3: Class Imbalance Handling",
    "Key Findings & Conclusions"
])

# Slide 3: Introduction
add_content_slide(prs, "Introduction", [
    "Financial fraud costs billions annually worldwide",
    "Traditional rule-based systems cannot adapt to evolving fraud patterns",
    "Machine learning offers flexible pattern detection from data",
    "Challenge: Extreme class imbalance (0.13% fraud rate)",
    "Challenge: Avoiding data leakage for realistic results"
])

# Slide 4: Research Questions
add_content_slide(prs, "Research Questions", [
    "RQ1: Which ML algorithms perform best for fraud detection in imbalanced data?",
    "",
    "RQ2: Which transaction features are most predictive of fraudulent activity?",
    "",
    "RQ3: How can class imbalance be handled to minimize false negatives while maintaining acceptable precision?",
    "",
    "Critical Focus: Using only pre-transaction features to avoid data leakage"
])

# Slide 5: Section - Dataset
add_section_slide(prs, "Dataset Overview")

# Slide 6: Dataset Stats
add_content_slide(prs, "PaySim Dataset", [
    "6,362,620 total transactions",
    "8,213 fraudulent cases (0.13%)",
    "5 transaction types:",
    "  CASH_IN, CASH_OUT, DEBIT,",
    "  PAYMENT, TRANSFER",
    "Fraud only in TRANSFER & CASH_OUT",
    "Class imbalance ratio: 773:1"
], f"{OUTPUT_DIR}/fig01_class_distribution.png", Inches(6))

# Slide 7: Fraud by Type
add_content_slide(prs, "Fraud by Transaction Type",
    image_path=f"{OUTPUT_DIR}/fig02_fraud_by_transaction_type.png",
    image_width=Inches(11))

# Slide 8: Amount Distribution
add_content_slide(prs, "Transaction Amount Distribution", [
    "Fraudulent transactions have:",
    "  - Larger average amounts",
    "  - Different distribution pattern",
    "",
    "Average Amount:",
    "  Legitimate: $178,197",
    "  Fraudulent: $1,467,967",
    "",
    "Frauds are 8.2x larger on average"
], f"{OUTPUT_DIR}/fig03_amount_distribution.png", Inches(6.5))

# Slide 9: Section - Data Leakage
add_section_slide(prs, "Critical Finding:\nData Leakage Prevention")

# Slide 10: Data Leakage Discovery
add_content_slide(prs, "Data Leakage Discovery", [
    "Initial models achieved >99% accuracy across all metrics",
    "",
    "Investigation revealed: full_drain feature had +0.987 correlation with fraud",
    "  - 100% of full_drain=1 transactions were fraudulent",
    "  - 97.5% of all frauds had full_drain=1",
    "",
    "Problem: full_drain requires POST-transaction data",
    "Real systems must predict BEFORE transaction completes",
    "",
    "This feature encodes the fraud definition itself!"
])

# Slide 11: Resolution
add_content_slide(prs, "Resolution: Pre-Transaction Features Only", [
    "Safe Features (10 total):",
    "  - step, hour_of_day, day (temporal)",
    "  - type_encoded (transaction type)",
    "  - amount (transaction amount)",
    "  - oldbalanceOrg, oldbalanceDest (pre-transaction balances)",
    "  - amount_to_orig_balance, amount_to_dest_balance (ratios)",
    "  - dest_zero_balance_before (recipient zero balance flag)",
    "",
    "Excluded Features (Data Leakage):",
    "  - full_drain, newbalanceOrig, newbalanceDest",
    "  - balance_change_orig, balance_change_dest, orig_balance_error"
])

# Slide 12: Section - RQ1
add_section_slide(prs, "RQ1: Algorithm Comparison")

# Slide 13: RQ1 Methods
add_content_slide(prs, "Algorithms Evaluated", [
    "Logistic Regression",
    "  - Linear baseline model",
    "",
    "Decision Tree",
    "  - Simple tree-based classifier",
    "",
    "Random Forest",
    "  - Ensemble of decision trees",
    "",
    "Gradient Boosting",
    "  - Sequential boosting approach",
    "",
    "XGBoost",
    "  - Optimized gradient boosting"
])

# Slide 14: RQ1 Results
add_content_slide(prs, "Algorithm Performance (Safe Features)", [
    "XGBoost: Best performer",
    "  AUC: 0.9977, Recall: 87.2%",
    "",
    "Gradient Boosting:",
    "  AUC: 0.9861, Recall: 61.0%",
    "",
    "Random Forest:",
    "  AUC: 0.9847, Recall: 52.8%",
    "",
    "Decision Tree:",
    "  AUC: 0.9554, Recall: 53.7%",
    "",
    "Logistic Regression:",
    "  AUC: 0.8258, Recall: 4.3%"
], f"{OUTPUT_DIR}/fig04_rq1_algorithm_comparison.png", Inches(6.5))

# Slide 15: ROC Curves
add_content_slide(prs, "ROC Curves Comparison",
    image_path=f"{OUTPUT_DIR}/fig05_rq1_roc_curves.png",
    image_width=Inches(9))

# Slide 16: PR Curves
add_content_slide(prs, "Precision-Recall Curves",
    image_path=f"{OUTPUT_DIR}/fig06_rq1_precision_recall_curves.png",
    image_width=Inches(9))

# Slide 17: RQ1 Conclusion
add_content_slide(prs, "RQ1 Key Findings", [
    "XGBoost significantly outperforms all other algorithms",
    "",
    "Tree-based ensemble methods dominate:",
    "  - Can capture non-linear fraud patterns",
    "  - Handle feature interactions effectively",
    "",
    "Linear models (Logistic Regression) fail:",
    "  - Only 4.3% recall - nearly useless",
    "  - Fraud patterns are highly non-linear",
    "",
    "Recommendation: XGBoost for production deployment"
])

# Slide 18: Section - RQ2
add_section_slide(prs, "RQ2: Feature Importance")

# Slide 19: Feature Importance
add_content_slide(prs, "Feature Importance Analysis", [
    "Top Features:",
    "",
    "1. amount_to_orig_balance: 97.7%",
    "   Ratio of amount to sender balance",
    "",
    "2. amount: 1.7%",
    "   Transaction amount",
    "",
    "3. type_encoded: 0.3%",
    "   Transaction type",
    "",
    "Dominant predictor: amount ratio"
], f"{OUTPUT_DIR}/fig07_rq2_feature_importance.png", Inches(6.5))

# Slide 20: Feature Correlation
add_content_slide(prs, "Feature Correlation with Fraud",
    image_path=f"{OUTPUT_DIR}/fig08_rq2_feature_correlation.png",
    image_width=Inches(10))

# Slide 21: Feature Distributions
add_content_slide(prs, "Top Feature Distributions by Class",
    image_path=f"{OUTPUT_DIR}/fig09_rq2_top_feature_distributions.png",
    image_width=Inches(11))

# Slide 22: RQ2 Conclusion
add_content_slide(prs, "RQ2 Key Findings", [
    "amount_to_orig_balance dominates (97.7% importance)",
    "  - Fraudsters transfer large proportions of account balances",
    "  - High ratios indicate account draining attempts",
    "",
    "Transaction amount is secondary indicator",
    "  - Fraudulent transactions are ~8x larger on average",
    "",
    "Transaction type provides context",
    "  - TRANSFER and CASH_OUT are fraud-associated types",
    "",
    "Key Insight: Amount ratios are more predictive than absolute amounts"
])

# Slide 23: Section - RQ3
add_section_slide(prs, "RQ3: Class Imbalance Handling")

# Slide 24: Imbalance Challenge
add_content_slide(prs, "The Class Imbalance Challenge", [
    "Dataset Distribution:",
    "  - Legitimate: 6,354,407 (99.87%)",
    "  - Fraudulent: 8,213 (0.13%)",
    "  - Imbalance ratio: 773:1",
    "",
    "Without handling, models:",
    "  - Predict everything as legitimate",
    "  - Miss most fraudulent transactions",
    "  - High accuracy but useless for fraud detection"
])

# Slide 25: Imbalance Methods
add_content_slide(prs, "Imbalance Handling Methods Evaluated", [
    "1. Baseline (no handling)",
    "   Standard XGBoost",
    "",
    "2. Random Undersampling",
    "   Reduce majority class",
    "",
    "3. SMOTE (Synthetic Oversampling)",
    "   Generate synthetic minority samples",
    "",
    "4. Class Weighting (scale_pos_weight)",
    "   Weight fraud samples higher",
    "",
    "5. Combined Approach",
    "   Undersampling + weighting"
], f"{OUTPUT_DIR}/fig10_rq3_imbalance_comparison.png", Inches(6.5))

# Slide 26: Precision-Recall Tradeoff
add_content_slide(prs, "Precision vs Recall Trade-off",
    image_path=f"{OUTPUT_DIR}/fig11_rq3_precision_recall_tradeoff.png",
    image_width=Inches(10))

# Slide 27: Confusion Matrices
add_content_slide(prs, "Confusion Matrices by Method",
    image_path=f"{OUTPUT_DIR}/fig12_rq3_confusion_matrices.png",
    image_width=Inches(12))

# Slide 28: RQ3 ROC Curves
add_content_slide(prs, "ROC Curves - Imbalance Methods",
    image_path=f"{OUTPUT_DIR}/fig13_rq3_roc_curves.png",
    image_width=Inches(9))

# Slide 29: Threshold Analysis
add_content_slide(prs, "Threshold Analysis",
    image_path=f"{OUTPUT_DIR}/fig15_rq3_threshold_analysis.png",
    image_width=Inches(11))

# Slide 30: RQ3 Results
add_content_slide(prs, "RQ3 Results Summary", [
    "Method               Recall    Precision   Missed Frauds",
    "-------------------------------------------------------",
    "Baseline XGBoost     87.2%     94.5%       315",
    "Undersampled         99.2%     24.4%       20",
    "SMOTE                98.8%     55.3%       29",
    "Weighted XGBoost     99.3%     60.3%       18",
    "Combined XGBoost     99.4%     36.9%       16",
    "",
    "Best Recall: Combined XGBoost (99.4%)",
    "Best Balance: Weighted XGBoost (99.3% recall, 60.3% precision)"
])

# Slide 31: Section - Conclusions
add_section_slide(prs, "Conclusions")

# Slide 32: Key Findings
add_content_slide(prs, "Key Findings", [
    "RQ1: XGBoost achieves best performance (AUC 0.9977)",
    "      Tree-based ensembles outperform linear models",
    "",
    "RQ2: amount_to_orig_balance is dominant predictor (97.7%)",
    "      Fraudsters drain large proportions of account balances",
    "",
    "RQ3: Class weighting provides best precision-recall balance",
    "      99.3% recall with 60.3% precision using Weighted XGBoost",
    "",
    "Critical: Data leakage identification ensures realistic expectations"
])

# Slide 33: Final Summary
add_content_slide(prs, "Summary Visualization",
    image_path=f"{OUTPUT_DIR}/fig16_final_summary.png",
    image_width=Inches(11))

# Slide 34: Practical Recommendations
add_content_slide(prs, "Practical Recommendations", [
    "For Production Deployment:",
    "",
    "1. Use XGBoost with class weighting (scale_pos_weight)",
    "",
    "2. Monitor amount_to_orig_balance as primary risk indicator",
    "",
    "3. Lower decision threshold (0.3-0.4) to catch more fraud",
    "   Accept higher false positives to minimize missed frauds",
    "",
    "4. Focus on TRANSFER and CASH_OUT transaction types",
    "",
    "5. Real-time monitoring of pre-transaction features only"
])

# Slide 35: Future Work
add_content_slide(prs, "Limitations & Future Work", [
    "Limitations:",
    "  - Synthetic dataset may not fully reflect real fraud patterns",
    "  - Limited feature engineering due to data leakage constraints",
    "  - No temporal pattern analysis across user history",
    "",
    "Future Directions:",
    "  - Real-world dataset validation",
    "  - Deep learning approaches (LSTM for sequences)",
    "  - Graph neural networks for transaction networks",
    "  - Online learning for adapting to new fraud patterns",
    "  - Cost-sensitive learning based on transaction amounts"
])

# Slide 36: Thank You
add_title_slide(prs, "Thank You", "Questions?")

# Save presentation
output_path = 'D:/Final Year/App Domains 3/datasets/Fraud_Detection_Presentation.pptx'
prs.save(output_path)

print("="*70)
print("PRESENTATION CREATED SUCCESSFULLY")
print("="*70)
print(f"\nSaved to: {output_path}")
print(f"\nTotal slides: {len(prs.slides)}")
print("\nSlide Overview:")
print("  1. Title Slide")
print("  2. Agenda")
print("  3-4. Introduction & Research Questions")
print("  5-8. Dataset Overview (with visualizations)")
print("  9-11. Data Leakage Discovery & Resolution")
print("  12-17. RQ1: Algorithm Comparison (4 visualizations)")
print("  18-22. RQ2: Feature Importance (3 visualizations)")
print("  23-30. RQ3: Imbalance Handling (5 visualizations)")
print("  31-35. Conclusions & Recommendations")
print("  36. Thank You")
print("\nVisualizations used: 13 figures from outputs/")
print("="*70)
